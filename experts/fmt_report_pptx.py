$extens("include.py")
include("import pptx", ["extella-pip install python-pptx"])

def fmt_report_pptx(input_path: str = "", records_json: str = "", spec_json: str = "",
                    output_path: str = "", brand_name: str = "", accent: str = "",
                    footer_note: str = "", sections_json: str = "") -> dict:
    """Оформитель отчётов в виде презентации. Тот же контракт, что у PDF и Word.

    Презентация — НЕ отчёт, разрезанный на слайды. Это другой жанр: её показывают людям,
    поэтому на слайде одна мысль, а не таблица целиком. Отсюда структура:
      титул → главное число → по слайду на разрез (топ-8) → закрывающий слайд.
    Выводов система НЕ выдумывает: она показывает числа, толкует их человек.

    Диаграммы РОДНЫЕ для PowerPoint (не картинки) — их можно править и пересчитывать,
    как таблицы в Word. Это и есть ценность формата: клиент доработает под свою встречу.

    Документ принадлежит КЛИЕНТУ: знаков Extella нет, в титуле имя владельца процесса.
    """
    import json
    import os
    from pathlib import Path
    from datetime import datetime, timezone
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    ready = []
    if sections_json and not sections_json.startswith("{{"):
        try:
            ready = [s for s in (json.loads(sections_json) or []) if isinstance(s, dict) and s.get("items")]
        except Exception:
            ready = []

    recs = []
    if records_json and not records_json.startswith("{{"):
        try:
            recs = json.loads(records_json)
        except Exception:
            return {"status": "error", "message": "records_json не разобрался"}
    elif input_path and not input_path.startswith("{{") and Path(input_path).exists():
        try:
            data = json.loads(Path(input_path).read_text(encoding="utf-8"))
        except Exception:
            return {"status": "error", "message": "вход не читается как JSON: " + str(input_path)[:120]}
        recs = data if isinstance(data, list) else (data.get("records") or data.get("rows") or [])
    recs = [r for r in (recs or []) if isinstance(r, dict)]
    if not ready and not recs:
        return {"status": "error", "message": "нет данных для презентации (ни records_json, ни input_path, ни sections_json)"}

    spec = {}
    if spec_json and not spec_json.startswith("{{"):
        try:
            spec = json.loads(spec_json) or {}
        except Exception:
            spec = {}
    style = spec.get("style") or {}
    accent = accent or style.get("accent") or "#2F6B66"
    brand_name = brand_name or style.get("brand_name") or ""
    footer_note = footer_note or style.get("footer") or ""
    cols = list(recs[0].keys()) if recs else []

    sections = list(ready)
    if not sections:
        views = [v for v in (spec.get("views") or []) if isinstance(v, dict) and v.get("group_by") in cols]
        if not views:
            cand = []
            for c in cols:
                vals = [str(r.get(c, "")).strip() for r in recs if str(r.get(c, "")).strip()]
                uniq = len(set(vals))
                if vals and 1 < uniq <= max(2, min(12, len(recs))) and uniq < len(vals):
                    cand.append((uniq, c))
            views = [{"group_by": c, "title": "По полю «" + c + "»"} for _, c in sorted(cand)[:3]]
        for v in views[:4]:
            g = v["group_by"]
            agg = {}
            for r in recs:
                k = str(r.get(g, "")).strip()
                agg[k or "не заполнено"] = agg.get(k or "не заполнено", 0) + 1
            named = {k: n for k, n in agg.items() if k != "не заполнено"}
            items = dict(sorted(named.items(), key=lambda x: -x[1])[:12])
            if agg.get("не заполнено"):
                items["не заполнено"] = agg["не заполнено"]
            sections.append({"title": v.get("title") or g, "items": items})

    hl = spec.get("headline") or {}
    if ready and not recs:
        _t = spec.get("total")
        hv = _t if _t is not None else sum(sum(s["items"].values()) for s in ready)
    else:
        hv = len(recs)
    hlabel = str(hl.get("label") or "записей в отчёте").replace("\n", " ")

    def _rgb(h):
        h = str(h or "").lstrip("#")
        if len(h) != 6:
            h = "2F6B66"
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except Exception:
            return RGBColor(0x2F, 0x6B, 0x66)

    A = _rgb(accent)
    INK, SILVER = RGBColor(0x1A, 0x1A, 0x1A), RGBColor(0x8C, 0x8C, 0x8C)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)   # 16:9 — как показывают сегодня
    BLANK = prs.slide_layouts[6]

    def _text(slide, txt, left, top, width, height, size, bold=False, color=INK):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(txt)
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Arial"
        return box

    # ── 1. Титул ──
    s1 = prs.slides.add_slide(BLANK)
    _text(s1, brand_name or " ", Inches(0.9), Inches(0.7), Inches(11), Inches(0.5), 14, False, SILVER)
    _text(s1, spec.get("report") or "Отчёт", Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.6), 40, True, INK)
    if spec.get("subtitle"):
        _text(s1, spec["subtitle"], Inches(0.9), Inches(4.0), Inches(11), Inches(1.0), 16, False, SILVER)
    _text(s1, datetime.now(timezone.utc).strftime("%d.%m.%Y"), Inches(0.9), Inches(6.3), Inches(5), Inches(0.5), 12, False, SILVER)

    # ── 2. Главное число: одна мысль на слайд ──
    s2 = prs.slides.add_slide(BLANK)
    _text(s2, str(hv), Inches(0.9), Inches(2.2), Inches(6), Inches(2.0), 96, True, A)
    _text(s2, hlabel, Inches(0.9), Inches(4.3), Inches(8), Inches(0.8), 20, False, INK)

    # ── 3. По слайду на разрез, диаграмма РОДНАЯ (правится в PowerPoint) ──
    for sec in sections[:6]:
        sl = prs.slides.add_slide(BLANK)
        _text(sl, sec["title"], Inches(0.9), Inches(0.6), Inches(11), Inches(0.8), 24, True, INK)
        items = list(sec["items"].items())[:8]   # больше восьми на слайде не читается
        cd = CategoryChartData()
        cd.categories = [str(k)[:34] for k, _ in items]
        cd.add_series("Количество", tuple(float(v) for _, v in items))
        gf = sl.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.9), Inches(1.6),
                                 Inches(11.5), Inches(5.0), cd)
        ch = gf.chart
        ch.has_legend = False
        try:
            ser = ch.plots[0].series[0]
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = A
            ch.plots[0].has_data_labels = True
            ch.plots[0].data_labels.font.size = Pt(11)
        except Exception:
            pass   # оформление диаграммы не обязано ронять презентацию
        if len(sec["items"]) > 8:
            _text(sl, "показаны 8 крупнейших из " + str(len(sec["items"])),
                  Inches(0.9), Inches(6.7), Inches(7), Inches(0.4), 10, False, SILVER)

    if footer_note:
        _text(prs.slides[0], footer_note, Inches(0.9), Inches(6.8), Inches(11), Inches(0.4), 10, False, SILVER)

    out = output_path or ("/tmp/report_" + datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S") + ".pptx")
    try:
        prs.save(out)
    except Exception as e:
        return {"status": "error", "message": "презентация не сохранилась: " + str(e)[:140]}
    return {"status": "success", "path": out, "bytes": os.path.getsize(out),
            "slides": len(prs.slides._sldIdLst), "records": len(recs),
            "preview": bool(ready and not recs), "sections": [s["title"] for s in sections]}
